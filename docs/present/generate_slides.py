"""
FinBridge Presentation Generator
Run: pip install python-pptx && python generate_slides.py
Output: FinBridge_Presentation.pptx (same folder)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color palette ─────────────────────────────────────────────────────────────
INDIGO       = RGBColor(79, 70, 229)
INDIGO_LIGHT = RGBColor(238, 242, 255)
INDIGO_DARK  = RGBColor(55, 48, 163)
WHITE        = RGBColor(255, 255, 255)
DARK         = RGBColor(17, 24, 39)
GRAY         = RGBColor(107, 114, 128)
GRAY_LIGHT   = RGBColor(249, 250, 251)
AMBER        = RGBColor(217, 119, 6)
GREEN        = RGBColor(22, 163, 74)
RED          = RGBColor(220, 38, 38)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=Pt(16), color=DARK, bullet="•", spacing=Inches(0.05)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def accent_bar(slide, color=INDIGO, height=Inches(0.12)):
    add_rect(slide, 0, 0, SLIDE_W, height, color)


def content_header(slide, title, subtitle=None):
    accent_bar(slide)
    add_text(slide, title,
             Inches(0.5), Inches(0.22), Inches(12), Inches(0.55),
             font_size=Pt(30), bold=True, color=INDIGO)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
                 font_size=Pt(15), color=GRAY)


def add_tag_box(slide, text, left, top, width, height, bg=INDIGO_LIGHT, fg=INDIGO):
    add_rect(slide, left, top, width, height, bg)
    add_text(slide, text, left + Inches(0.15), top + Inches(0.08),
             width - Inches(0.3), height - Inches(0.15),
             font_size=Pt(13), bold=True, color=fg, align=PP_ALIGN.CENTER)


# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, INDIGO)

    # Diagonal accent block
    add_rect(slide, Inches(8.5), 0, Inches(5), SLIDE_H, INDIGO_DARK)

    # Main title
    add_text(slide, "FinBridge",
             Inches(0.7), Inches(1.8), Inches(8), Inches(1.4),
             font_size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_text(slide, "AI-Powered Financial Document Management Platform",
             Inches(0.7), Inches(3.3), Inches(8.2), Inches(0.9),
             font_size=Pt(22), color=RGBColor(199, 210, 254), align=PP_ALIGN.LEFT)

    # Tagline
    add_text(slide, "Replace WhatsApp invoicing with structured AI-powered workflows",
             Inches(0.7), Inches(4.15), Inches(8.2), Inches(0.6),
             font_size=Pt(15), color=RGBColor(165, 180, 252), align=PP_ALIGN.LEFT)

    # Hackathon tag
    add_rect(slide, Inches(0.7), Inches(5.0), Inches(2.2), Inches(0.45),
             RGBColor(55, 48, 163))
    add_text(slide, "  Hackathon 2025",
             Inches(0.7), Inches(5.0), Inches(2.2), Inches(0.45),
             font_size=Pt(13), bold=True, color=WHITE)

    # Tech badges
    badges = ["Claude claude-sonnet-4-6", "FastAPI", "React + Vite", "PostgreSQL", "PWA"]
    x = Inches(0.7)
    for b in badges:
        w = Inches(len(b) * 0.115 + 0.3)
        add_rect(slide, x, Inches(5.7), w, Inches(0.38), RGBColor(67, 56, 202))
        add_text(slide, b, x + Inches(0.1), Inches(5.72), w - Inches(0.1), Inches(0.35),
                 font_size=Pt(12), color=WHITE)
        x += w + Inches(0.15)


# ── Slide 2 — Problem ─────────────────────────────────────────────────────────

def slide_problem(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    content_header(slide, "The Problem", "How financial documents flow today — and why it breaks")

    # Left: current state box
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5),
             RGBColor(254, 242, 242))
    add_text(slide, "Today's Reality",
             Inches(0.65), Inches(1.6), Inches(5.5), Inches(0.45),
             font_size=Pt(16), bold=True, color=RED)

    current_items = [
        "Company WhatsApps invoice photo to accountant",
        "Accountant manually re-types every field",
        "No visibility — did they receive it?",
        "Mistakes from manual entry go undetected",
        "Duplicate invoices processed twice",
        "50-employee salary register = 50 manual rows",
        "Zero audit trail — no record of changes",
    ]
    add_bullet_list(slide, current_items,
                    Inches(0.65), Inches(2.15), Inches(5.6), Inches(4.5),
                    font_size=Pt(14), color=RGBColor(127, 29, 29), bullet="✗")

    # Right: impact box
    add_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5),
             RGBColor(239, 246, 255))
    add_text(slide, "The Cost",
             Inches(6.95), Inches(1.6), Inches(5.6), Inches(0.45),
             font_size=Pt(16), bold=True, color=INDIGO)

    impact_items = [
        "Hours lost per week on data entry per accountant",
        "Errors in financial records → compliance risk",
        "No structured handoff → files lost in threads",
        "No status visibility → constant follow-up calls",
        "Impossible to scale — linear with document volume",
    ]
    add_bullet_list(slide, impact_items,
                    Inches(6.95), Inches(2.15), Inches(5.7), Inches(4.5),
                    font_size=Pt(14), color=DARK, bullet="→")


# ── Slide 3 — Solution ────────────────────────────────────────────────────────

def slide_solution(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    content_header(slide, "Our Solution — FinBridge",
                   "Structured AI-powered handoff between companies and their accounting firm")

    # Flow steps
    steps = [
        ("📤", "Upload", "Company uploads PDF/image via portal or mobile camera"),
        ("🤖", "AI Extract", "Claude vision extracts all fields with confidence scoring"),
        ("👁", "Review", "Accountant reviews, corrects amber-flagged fields"),
        ("✅", "Accept", "Transaction accepted with payment head assigned"),
        ("🔔", "Notify", "Company sees real-time status update"),
    ]

    step_w = Inches(2.3)
    gap    = Inches(0.25)
    start_x = Inches(0.4)
    top = Inches(1.8)

    for i, (icon, title, desc) in enumerate(steps):
        x = start_x + i * (step_w + gap)

        # Box
        add_rect(slide, x, top, step_w, Inches(3.2), INDIGO_LIGHT)

        # Icon
        add_text(slide, icon, x, top + Inches(0.2), step_w, Inches(0.6),
                 font_size=Pt(28), align=PP_ALIGN.CENTER)

        # Step number
        add_text(slide, f"Step {i+1}",
                 x, top + Inches(0.85), step_w, Inches(0.3),
                 font_size=Pt(11), color=INDIGO, align=PP_ALIGN.CENTER, bold=True)

        # Title
        add_text(slide, title, x, top + Inches(1.15), step_w, Inches(0.45),
                 font_size=Pt(17), bold=True, color=DARK, align=PP_ALIGN.CENTER)

        # Description
        add_text(slide, desc, x + Inches(0.1), top + Inches(1.65),
                 step_w - Inches(0.2), Inches(1.3),
                 font_size=Pt(12), color=GRAY, align=PP_ALIGN.CENTER, wrap=True)

        # Arrow between steps
        if i < len(steps) - 1:
            add_text(slide, "→",
                     x + step_w, top + Inches(1.3), gap, Inches(0.5),
                     font_size=Pt(20), color=INDIGO, align=PP_ALIGN.CENTER, bold=True)

    # Value props bar
    add_rect(slide, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.7), INDIGO)

    props = [
        ("6", "Document Types Supported"),
        ("5", "User Roles with RBAC"),
        ("Real-time", "Notifications via SSE"),
        ("100%", "Audit Trail Coverage"),
        ("PWA", "Works on Mobile"),
    ]
    prop_w = Inches(2.5)
    for i, (val, label) in enumerate(props):
        px = Inches(0.4) + i * prop_w
        add_text(slide, val,
                 px, Inches(5.4), prop_w, Inches(0.6),
                 font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, label,
                 px, Inches(5.95), prop_w, Inches(0.5),
                 font_size=Pt(11), color=RGBColor(199, 210, 254), align=PP_ALIGN.CENTER)


# ── Slide 4 — Architecture ────────────────────────────────────────────────────

def slide_architecture(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    content_header(slide, "Technical Architecture",
                   "Production-grade design decisions — built to scale")

    # Left column: stack layers
    add_text(slide, "Stack Layers",
             Inches(0.5), Inches(1.55), Inches(5.8), Inches(0.4),
             font_size=Pt(14), bold=True, color=INDIGO)

    layers = [
        ("Frontend", "React 18 + Vite + TypeScript + Tailwind v4", INDIGO_LIGHT),
        ("State", "React Query (server) + Zustand (client)", RGBColor(240, 253, 244)),
        ("Backend", "Python 3.12 + FastAPI + Pydantic", RGBColor(255, 251, 235)),
        ("Database", "PostgreSQL 16 + SQLAlchemy 2.0 + Alembic", RGBColor(253, 242, 248)),
        ("AI", "Claude claude-sonnet-4-6 — vision + tool_use", INDIGO_LIGHT),
        ("Auth", "JWT (python-jose) + bcrypt — stateless RBAC", RGBColor(240, 253, 244)),
        ("Real-time", "Server-Sent Events — no extra infra", RGBColor(255, 251, 235)),
        ("PWA", "vite-plugin-pwa + Workbox service worker", RGBColor(253, 242, 248)),
    ]

    for i, (layer, desc, bg) in enumerate(layers):
        y = Inches(2.0) + i * Inches(0.63)
        add_rect(slide, Inches(0.5), y, Inches(6.0), Inches(0.58), bg)
        add_text(slide, layer,
                 Inches(0.6), y + Inches(0.06), Inches(1.3), Inches(0.45),
                 font_size=Pt(12), bold=True, color=INDIGO)
        add_text(slide, desc,
                 Inches(1.95), y + Inches(0.06), Inches(4.4), Inches(0.45),
                 font_size=Pt(12), color=DARK)

    # Right column: design principles
    add_text(slide, "Key Design Decisions",
             Inches(7.0), Inches(1.55), Inches(6.0), Inches(0.4),
             font_size=Pt(14), bold=True, color=INDIGO)

    decisions = [
        ("🔐 Tenant Isolation", "company_id injected from JWT server-side.\nClient never sends tenant context."),
        ("🗄 DB-Driven Doc Types", "Claude prompt + field schema stored in DB.\nNew type = 1 SQL INSERT, zero code."),
        ("🤖 tool_use not plain text", "Forces Claude to return strict JSON schema.\nNo parsing, no hallucinated field names."),
        ("📡 SSE not WebSocket", "Unidirectional push — simpler infra.\nBrowser-native, no extra dependencies."),
        ("🆔 UUID Primary Keys", "No sequential ID enumeration attacks.\nAll entities use uuid4."),
    ]

    for i, (title, desc) in enumerate(decisions):
        y = Inches(2.0) + i * Inches(0.98)
        add_rect(slide, Inches(7.0), y, Inches(5.8), Inches(0.88), GRAY_LIGHT)
        add_text(slide, title,
                 Inches(7.15), y + Inches(0.05), Inches(5.5), Inches(0.35),
                 font_size=Pt(13), bold=True, color=DARK)
        add_text(slide, desc,
                 Inches(7.15), y + Inches(0.38), Inches(5.5), Inches(0.48),
                 font_size=Pt(11), color=GRAY, wrap=True)


# ── Slide 5 — AI Capabilities ─────────────────────────────────────────────────

def slide_ai(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    content_header(slide, "Two Layers of AI Intelligence",
                   "Claude powers both document extraction and financial insights")

    # Left box: Document Extraction
    add_rect(slide, Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.6), INDIGO_LIGHT)
    add_text(slide, "📄  Document Extraction",
             Inches(0.55), Inches(1.65), Inches(5.7), Inches(0.55),
             font_size=Pt(18), bold=True, color=INDIGO)
    add_text(slide, "Claude claude-sonnet-4-6 Vision + tool_use",
             Inches(0.55), Inches(2.2), Inches(5.7), Inches(0.4),
             font_size=Pt(13), color=GRAY, bold=True)

    extraction_points = [
        "Prompt template + field schema stored in DB per document type",
        "tool_use forces deterministic JSON — no parsing needed",
        "confidence_score + low_confidence_fields per transaction",
        "Native PDF support (no image conversion needed)",
        "6 document types: Invoice, Salary, Bank Statement, and more",
        "Adding a new type = 1 SQL INSERT, zero code changes",
        "Full raw_ai_output stored for audit and debugging",
    ]
    add_bullet_list(slide, extraction_points,
                    Inches(0.55), Inches(2.65), Inches(5.7), Inches(4.0),
                    font_size=Pt(13), color=DARK)

    # Right box: AI Insights Panel
    add_rect(slide, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.6), RGBColor(240, 253, 244))
    add_text(slide, "✨  AI Insights Panel",
             Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.55),
             font_size=Pt(18), bold=True, color=GREEN)
    add_text(slide, "Pre-built queries + Claude natural language explanation",
             Inches(7.05), Inches(2.2), Inches(5.7), Inches(0.4),
             font_size=Pt(13), color=GRAY, bold=True)

    insights_points = [
        "7 pre-built financial queries — no Text-to-SQL hallucination risk",
        "company_id injected server-side — tenant isolation guaranteed",
        "Claude explains results in 2-3 plain English sentences",
        "Pending > ₹1 Lakh  |  Spend by Head  |  Top Vendors",
        "Low Confidence  |  Rejected + Reasons  |  Monthly Trend",
        "Floats on every portal — company, accountant, firm admin",
        "React Portal ensures widget renders outside flex containers",
    ]
    add_bullet_list(slide, insights_points,
                    Inches(7.05), Inches(2.65), Inches(5.7), Inches(4.0),
                    font_size=Pt(13), color=DARK)

    # vs label in the middle
    add_text(slide, "VS\nText-to-SQL",
             Inches(6.05), Inches(3.2), Inches(0.9), Inches(1.0),
             font_size=Pt(11), bold=True, color=GRAY, align=PP_ALIGN.CENTER)


# ── Slide 6 — Demo Flow ───────────────────────────────────────────────────────

def slide_demo(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    content_header(slide, "End-to-End Demo",
                   "Full workflow visible in under 90 seconds")

    steps = [
        ("1", "admin@techcorp.com", "Upload PDF invoice via drag-and-drop or mobile camera", INDIGO),
        ("2", "System (Claude)", "Status: pending → processing → extracted in real time", AMBER),
        ("3", "accountant1@apex.com", "Bell notification — open Review Queue, see amber-flagged fields", RGBColor(234, 88, 12)),
        ("4", "accountant1@apex.com", "Correct fields, assign payment head, click Accept", GREEN),
        ("5", "admin@techcorp.com", "Bell notification: 'Transaction accepted' — status updated", INDIGO),
        ("6", "firm@apexaccounting.com", "Audit Trail shows every action with user + timestamp", GRAY),
        ("7", "Any user", "Click ✨ Ask AI → 'Pending > ₹1 Lakh' → Claude explains live DB results", GREEN),
    ]

    for i, (num, role, action, color) in enumerate(steps):
        y = Inches(1.5) + i * Inches(0.76)

        # Number circle
        add_rect(slide, Inches(0.4), y + Inches(0.08), Inches(0.45), Inches(0.45), color)
        add_text(slide, num,
                 Inches(0.4), y + Inches(0.05), Inches(0.45), Inches(0.48),
                 font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Role tag
        tag_w = Inches(2.6)
        add_rect(slide, Inches(1.0), y + Inches(0.08), tag_w, Inches(0.42),
                 RGBColor(243, 244, 246))
        add_text(slide, role,
                 Inches(1.05), y + Inches(0.1), tag_w - Inches(0.1), Inches(0.38),
                 font_size=Pt(11), bold=True, color=color)

        # Action
        add_text(slide, action,
                 Inches(3.75), y + Inches(0.1), Inches(9.2), Inches(0.45),
                 font_size=Pt(13), color=DARK)

    # Credentials footer
    add_rect(slide, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.55),
             RGBColor(243, 244, 246))
    creds = "Demo credentials — TechCorp Admin: admin@techcorp.com / Tech@1234   |   " \
            "Accountant: accountant1@apex.com / Acc@1234   |   Firm Admin: firm@apexaccounting.com / Firm@1234"
    add_text(slide, creds,
             Inches(0.55), Inches(6.87), Inches(12.2), Inches(0.48),
             font_size=Pt(10), color=GRAY)


# ── Slide 7 — What's Next ─────────────────────────────────────────────────────

def slide_next(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    content_header(slide, "What We'd Build Next",
                   "The foundation is production-ready — here's the roadmap")

    roadmap = [
        ("☁️", "S3 File Storage",
         "Replace local uploads/ with AWS S3. Already abstracted behind UPLOAD_DIR — a config change.",
         "Scale"),
        ("🏛", "GST Validation API",
         "Auto-verify GSTIN against government API. Flag mismatches before accountant review.",
         "Compliance"),
        ("💬", "WhatsApp Bot",
         "Company sends bill photo to a WhatsApp number → auto-extracted and loaded into FinBridge. Removes portal upload entirely.",
         "Reach"),
        ("💱", "Multi-Currency",
         "Handle USD/EUR invoices. Store original currency, display INR equivalent via live exchange rates.",
         "Global"),
        ("🧠", "ML Payment Head Suggestions",
         "Train on accepted transactions. Auto-suggest payment heads — learns from every accountant correction.",
         "AI"),
    ]

    card_w = Inches(2.3)
    gap    = Inches(0.3)
    start_x = Inches(0.4)

    for i, (icon, title, desc, tag) in enumerate(roadmap):
        x = start_x + i * (card_w + gap)

        add_rect(slide, x, Inches(1.6), card_w, Inches(4.8), GRAY_LIGHT)

        # Tag badge
        add_rect(slide, x + card_w - Inches(0.9), Inches(1.65), Inches(0.85), Inches(0.32), INDIGO)
        add_text(slide, tag,
                 x + card_w - Inches(0.9), Inches(1.65), Inches(0.85), Inches(0.32),
                 font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Icon
        add_text(slide, icon,
                 x, Inches(2.05), card_w, Inches(0.6),
                 font_size=Pt(30), align=PP_ALIGN.CENTER)

        # Title
        add_text(slide, title,
                 x + Inches(0.1), Inches(2.75), card_w - Inches(0.2), Inches(0.55),
                 font_size=Pt(15), bold=True, color=DARK, align=PP_ALIGN.CENTER)

        # Description
        add_text(slide, desc,
                 x + Inches(0.12), Inches(3.38), card_w - Inches(0.24), Inches(2.7),
                 font_size=Pt(12), color=GRAY, align=PP_ALIGN.CENTER, wrap=True)

    # Closing line
    add_rect(slide, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.72), INDIGO)
    add_text(slide,
             "Built in a hackathon — designed for production. Multi-tenant RBAC, real-time SSE, "
             "PWA-ready, audit-complete, AI at the core.",
             Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.65),
             font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slide 1 — Title...")
    slide_title(prs)

    print("Building slide 2 — Problem...")
    slide_problem(prs)

    print("Building slide 3 — Solution...")
    slide_solution(prs)

    print("Building slide 4 — Architecture...")
    slide_architecture(prs)

    print("Building slide 5 — AI Capabilities...")
    slide_ai(prs)

    print("Building slide 6 — Demo Flow...")
    slide_demo(prs)

    print("Building slide 7 — What's Next...")
    slide_next(prs)

    out = os.path.join(os.path.dirname(__file__), "FinBridge_Presentation.pptx")
    prs.save(out)
    print(f"\nSaved: {out}")
    print("   Open in PowerPoint or Google Slides (File > Import > Upload)")


if __name__ == "__main__":
    main()
